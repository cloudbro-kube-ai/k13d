#!/usr/bin/env python3
"""Capture each slide from HTML and create PPTX with screenshots."""

import os
import glob
from pathlib import Path
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Emu

PRESENTATION_DIR = Path(__file__).parent
HTML_FILE = PRESENTATION_DIR / "index.html"
OUTPUT_PPTX = PRESENTATION_DIR / "k13d_리뉴얼_발표자료.pptx"
SCREENSHOT_DIR = PRESENTATION_DIR / "screenshots"

# 16:9 slide dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def capture_slides():
    """Capture each slide as PNG using Playwright."""
    SCREENSHOT_DIR.mkdir(exist_ok=True)
    
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": 1920, "height": 1080})
        
        # Load HTML file
        page.goto(f"file://{HTML_FILE.resolve()}")
        page.wait_for_timeout(1000)
        
        # Get total slide count
        total_slides = page.evaluate("document.querySelectorAll('.slide').length")
        print(f"Total slides: {total_slides}")
        
        # Hide navigation hint and PDF button
        page.evaluate("""
            document.querySelector('.nav-hint').style.display = 'none';
            document.querySelector('.pdf-download-btn').style.display = 'none';
        """)
        
        for i in range(total_slides):
            # Show only current slide
            page.evaluate(f"""
                const slides = document.querySelectorAll('.slide');
                slides.forEach((s, idx) => {{
                    s.style.display = idx === {i} ? 'flex' : 'none';
                }});
            """)
            page.wait_for_timeout(300)
            
            # Capture screenshot
            screenshot_path = SCREENSHOT_DIR / f"slide_{i+1:02d}.png"
            page.screenshot(path=str(screenshot_path), full_page=False)
            print(f"Captured: slide_{i+1:02d}.png")
        
        browser.close()
    
    return total_slides


def create_pptx(total_slides):
    """Create PPTX from captured screenshots."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # Use blank layout
    blank_layout = prs.slide_layouts[6]
    
    for i in range(1, total_slides + 1):
        screenshot_path = SCREENSHOT_DIR / f"slide_{i:02d}.png"
        if not screenshot_path.exists():
            print(f"Warning: {screenshot_path} not found, skipping")
            continue
        
        slide = prs.slides.add_slide(blank_layout)
        
        # Add screenshot as full-slide image
        slide.shapes.add_picture(
            str(screenshot_path),
            Emu(0),
            Emu(0),
            SLIDE_WIDTH,
            SLIDE_HEIGHT
        )
        
        print(f"Added slide {i}")
    
    prs.save(str(OUTPUT_PPTX))
    print(f"\nPPTX saved: {OUTPUT_PPTX}")


if __name__ == "__main__":
    print("=== Capturing slides from HTML ===")
    total = capture_slides()
    
    print("\n=== Creating PPTX ===")
    create_pptx(total)
    
    # Cleanup screenshots
    print("\n=== Cleanup ===")
    for f in SCREENSHOT_DIR.glob("*.png"):
        f.unlink()
    SCREENSHOT_DIR.rmdir()
    print("Screenshots cleaned up")
