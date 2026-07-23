#!/usr/bin/env python3
"""Generate k13d presentation PPTX from HTML content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT_RED = RGBColor(0xe9, 0x45, 0x60)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf8, 0xf9, 0xfa)
DARK_TEXT = RGBColor(0x1a, 0x1a, 0x2e)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)
BLUE = RGBColor(0x0f, 0x34, 0x60)
GREEN = RGBColor(0x16, 0xa0, 0x85)
PURPLE = RGBColor(0x9b, 0x59, 0xb6)
ORANGE = RGBColor(0xf3, 0x9c, 0x12)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide_number(slide, num, total):
    txBox = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY_TEXT
    p.alignment = PP_ALIGN.RIGHT

def add_shape(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_paragraph(tf, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p

TOTAL_SLIDES = 20

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK_BG

add_text_box(slide, Inches(2), Inches(1.5), Inches(9), Inches(1.2), "k13d", 72, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(2.8), Inches(9), Inches(0.6), "Kubernetes AI Dashboard", 28, False, RGBColor(0xff, 0xff, 0xff), PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.5), "Reboot Project - Feature Branch Review", 16, False, RGBColor(0xaa, 0xaa, 0xaa), PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(4.3), Inches(9), Inches(0.5), "2026.03 ~ 2026.06", 16, False, RGBColor(0xaa, 0xaa, 0xaa), PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.5), Inches(5), Inches(6), Inches(0.5), "37 Commits | 73 Files Changed | +9,169 Lines", 14, False, RGBColor(0xcc, 0xcc, 0xcc), PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.5), "github.com/cloudbro-kube-ai/k13d", 13, False, RGBColor(0x88, 0x88, 0x88), PP_ALIGN.CENTER)
add_slide_number(slide, 1, TOTAL_SLIDES)

# Slide 2: About
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6), "K13D: 지능형 DevOps 플랫폼", 36, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

tf = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), "k + 13 letters + d = kubeaidashboard", 20, True, DARK_TEXT)
add_paragraph(tf, "K8s 운영의 새로운 패러다임", 24, True, DARK_TEXT, space_before=8)
add_paragraph(tf, "Dashboard + AI Assistant", 28, True, ACCENT_RED, space_before=4)

add_text_box(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(2),
    "k9s, Kubernetes Dashboard, kubectl-ai — K8s 도구들이 각각의 인터페이스에 갇혀 단편화된 불편함을 하나의 바이너리로 해결합니다.\n\nTUI와 Web UI를 동시에 제공하며, AI가 리소스 컨텍스트를 이해하고 실제 kubectl 명령어를 직접 실행하는 올인원 DevOps 플랫폼입니다.",
    14, False, GRAY_TEXT)

# Comparison boxes
comparisons = [
    ("vs k9s", "TUI 전용 → TUI + Web + AI"),
    ("vs K8s Dashboard", "Web 전용 → Web + TUI + AI"),
    ("vs kubectl-ai", "CLI 전용 → Dashboard + TUI"),
    ("vs Lens", "Desktop IDE → 단일 바이너리 + AI"),
]
for i, (vs, label) in enumerate(comparisons):
    row, col = divmod(i, 2)
    x = Inches(0.8 + col * 2.9)
    y = Inches(5 + row * 0.8)
    shape = add_shape(slide, x, y, Inches(2.7), Inches(0.7), RGBColor(0xff, 0xf0, 0xf2), ACCENT_RED)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = vs
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY_TEXT

# Right side mockup
add_shape(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(5.5), WHITE, RGBColor(0xe0, 0xe0, 0xe0))
add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.4), "k13d dashboard", 12, True, GRAY_TEXT)

pods = [
    ("nginx-pod", "default", "Running", GREEN),
    ("api-server", "prod", "Running", GREEN),
    ("redis-master", "prod", "CrashLoop", ACCENT_RED),
    ("gateway", "default", "Running", GREEN),
    ("worker-1", "data", "OOMKill", ORANGE),
]
for i, (name, ns, status, color) in enumerate(pods):
    y = Inches(2.2 + i * 0.45)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.3), y + Inches(0.08), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text_box(slide, Inches(7.6), y, Inches(2), Inches(0.35), name, 11, True, DARK_TEXT)
    add_text_box(slide, Inches(9.5), y, Inches(1), Inches(0.35), ns, 9, False, GRAY_TEXT)
    add_text_box(slide, Inches(10.5), y, Inches(1.5), Inches(0.35), status, 9, True, color)

add_slide_number(slide, 2, TOTAL_SLIDES)

# Slide 3: Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6), "k13D Architecture", 36, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.15), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.3), "USER INTERFACE", 10, True, GRAY_TEXT, PP_ALIGN.CENTER)

ui_boxes = [("TUI", "tview · Terminal"), ("Web UI", "HTTP · Browser"), ("CLI / MCP", "REPL · stdio")]
for i, (name, sub) in enumerate(ui_boxes):
    x = Inches(2 + i * 3.5)
    shape = add_shape(slide, x, Inches(1.9), Inches(2.5), Inches(0.8), WHITE, RGBColor(0xe0, 0xe0, 0xe0))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY_TEXT
    p2.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(0.8), Inches(3), Inches(11), Inches(0.3), "SHARED CORE", 10, True, GRAY_TEXT, PP_ALIGN.CENTER)

core_shape = add_shape(slide, Inches(1.5), Inches(3.4), Inches(10), Inches(0.9), DARK_BG)
core_items = [("AI Agent", "Providers · Tools · Safety"), ("K8s Client", "client-go · Watcher"), ("Config & DB", "Settings · Audit · Session"), ("Security", "Safety · Scanner")]
for i, (name, sub) in enumerate(core_items):
    x = Inches(1.8 + i * 2.5)
    add_text_box(slide, x, Inches(3.45), Inches(2), Inches(0.3), name, 12, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.75), Inches(2), Inches(0.3), sub, 9, False, RGBColor(0x88, 0x88, 0x88), PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.3), "EXTERNAL SERVICES", 10, True, GRAY_TEXT, PP_ALIGN.CENTER)

ext_items = [("Kubernetes API", "Cluster Resources", "client-go"), ("LLM Provider", "OpenAI / Ollama / ...", "REST API"), ("MCP Server", "Extensible Tools", "JSON-RPC"), ("SQLite", "Audit Log & Settings", "Embedded")]
for i, (name, sub, badge) in enumerate(ext_items):
    x = Inches(1.2 + i * 3)
    shape = add_shape(slide, x, Inches(5), Inches(2.5), Inches(1), WHITE, RGBColor(0xe0, 0xe0, 0xe0))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.size = Pt(9)
    p2.font.color.rgb = GRAY_TEXT
    p2.alignment = PP_ALIGN.CENTER

add_slide_number(slide, 3, TOTAL_SLIDES)

# Slide 4: Feature Map
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "K13D 기능 구성도", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

groups = [
    ("K8s Resource Views", "25+ types", BLUE, "Core: Pods, Deployments, Services, Nodes\nConfig: ConfigMaps, Secrets, PV/PVC\nWorkloads: ReplicaSets, DaemonSets, Jobs\nNetworking: Ingresses, NetworkPolicies"),
    ("AI Assistant", "7+ providers", ACCENT_RED, "NL Query: Natural language K8s 질의\nContext-Aware: YAML + Events + Logs\nTool Use: kubectl, bash, MCP\nSafety: 위험 명령어 승인 모달"),
    ("TUI Operations", "k9s parity", GREEN, "Navigation: Vim-style j/k/g/G\nCommand Mode: :pods, :deploy, :svc\nResource Actions: describe, YAML, edit\nFilter & Sort: / 필터, 정렬"),
    ("Web UI", "full dashboard", PURPLE, "Dashboard: 리소스 대시보드\nAI Chat: Web 환경 AI 비서\nMulti-Cluster: 여러 클러스터 전환\nAuth: Local/OIDC/LDAP 인증"),
    ("Monitoring", "deep insights", ORANGE, "Metrics: Pod/Node CPU/Memory\nPulse: 클러스터 상태 진단\nXRay: 리소스 계층 분석\nAnalysis: AI 기반 분석"),
    ("Operations", "enterprise", RGBColor(0xc0, 0x39, 0x2b), "GitOps: ArgoCD, Flux 연동\nBackup: Velero 백업/복구\nSelf-Healing: 자동 복구\nAudit: SQLite 감사 로그"),
]

for i, (title, badge, color, content) in enumerate(groups):
    row, col = divmod(i, 3)
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.2 + row * 3.1)
    
    shape = add_shape(slide, x, y, Inches(3.9), Inches(2.9), WHITE)
    shape.line.color.rgb = RGBColor(0xe0, 0xe0, 0xe0)
    shape.line.width = Pt(0.5)
    
    # Header
    header = add_shape(slide, x, y, Inches(3.9), Inches(0.5), color)
    tf = header.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{title}  [{badge}]"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # Content
    add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(3.6), Inches(2.2), content, 10, False, GRAY_TEXT)

add_slide_number(slide, 4, TOTAL_SLIDES)

# Slide 5: API Spec
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "K13D API Spec", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

api_groups = [
    ("System [4]", ["GET /api/health", "GET /api/version", "GET /api/features", "GET /api/metrics"]),
    ("Auth [12]", ["POST /api/auth/login", "POST /api/auth/logout", "GET /api/auth/status", "OIDC login/callback"]),
    ("K8s Resources [10]", ["GET /api/k8s/{resource}", "POST /api/k8s/apply", "GET /api/overview", "GET /api/pulse"]),
    ("Workload Ops [16]", ["PUT deployment/scale", "POST restart/pause", "GET deployment/history", "POST node/drain"]),
    ("AI & LLM [18]", ["POST /api/chat/agentic", "GET /api/sessions", "GET /api/models", "GET /api/mcp/servers"]),
    ("Cluster [12]", ["GET /api/contexts", "GET /api/pods/{name}/logs", "WebSocket /api/terminal", "GET /api/topology"]),
]

for i, (title, endpoints) in enumerate(api_groups):
    row, col = divmod(i, 3)
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.2 + row * 3.1)
    
    shape = add_shape(slide, x, y, Inches(3.9), Inches(2.9), WHITE)
    shape.line.color.rgb = RGBColor(0xe0, 0xe0, 0xe0)
    shape.line.width = Pt(0.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(8)
    
    for ep in endpoints:
        p2 = tf.add_paragraph()
        p2.text = ep
        p2.font.size = Pt(10)
        p2.font.color.rgb = GRAY_TEXT
        p2.font.name = "Monaco"
        p2.space_before = Pt(4)

add_slide_number(slide, 5, TOTAL_SLIDES)

# Slide 6: LLM Providers
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "k13d에서 사용가능한 LLM", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

# Table
rows = 10
cols = 7
table_shape = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.5))
table = table_shape.table

headers = ["Provider", "유형", "Tool Calling", "Streaming", "Multi-Model", "환경", "비고"]
providers = [
    ["OpenAI", "REST API", "Function Call", "SSE", "O", "Cloud", "GPT-4o, GPT-4, o-series"],
    ["Anthropic", "REST API", "Native Tool", "SSE", "O", "Cloud", "Claude 4, Claude 3.5 Sonnet"],
    ["Google Gemini", "REST API", "Function Call", "SSE", "O", "Cloud", "Gemini 2.5 Flash/Pro"],
    ["AWS Bedrock", "AWS SDK", "Tool Use", "SSE", "O", "Enterprise", "Claude, Llama 등 다수 모델"],
    ["Azure OpenAI", "REST API", "Function Call", "SSE", "O", "Enterprise", "Microsoft Azure 클라우드"],
    ["Ollama", "REST API", "Native Tool", "NDJSON", "O", "Local", "LLaMA, Mistral 등 로컬 실행"],
    ["LiteLLM", "Proxy", "Delegate", "SSE", "O", "Proxy", "100+ LLM 프록시 지원"],
    ["OpenRouter", "Proxy", "Function Call", "SSE", "O", "Proxy", "다수 모델 라우팅"],
    ["Upstage Solar", "REST API", "Function Call", "SSE", "O", "Cloud", "한국형 LLM (Solar Pro2)"],
]

for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BG

for row_idx, row_data in enumerate(providers):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_TEXT
        if row_idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xf5, 0xf5, 0xf5)

add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
    "Tip: 모든 Provider는 config.yaml에서 llm.provider 필드로 설정하며, Tool Calling은 kubectl, bash, MCP 도구를 AI가 직접 실행할 때 필요합니다.",
    10, False, GRAY_TEXT)

add_slide_number(slide, 6, TOTAL_SLIDES)

# Slide 7: AI Answering Flow
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "K13D가 답변을 내는 과정", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

flow_items = [
    ("사용자 질문", '"redis-master가 왜 CrashLoop인가요?"', DARK_BG),
    ("컨텍스트 수집", "YAML + Events + Logs", BLUE),
    ("LLM 호출", "Tool Call 요청", PURPLE),
]
for i, (title, sub, color) in enumerate(flow_items):
    x = Inches(0.8 + i * 4.2)
    shape = add_shape(slide, x, Inches(1.3), Inches(3.5), Inches(1), WHITE, color)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = sub
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY_TEXT
    p2.alignment = PP_ALIGN.CENTER

# Tool execution
tools_shape = add_shape(slide, Inches(0.5), Inches(2.7), Inches(12), Inches(1.2), WHITE, RGBColor(0xe0, 0xe0, 0xe0))
add_text_box(slide, Inches(0.7), Inches(2.8), Inches(3), Inches(0.3), "도구 실행", 12, True, DARK_TEXT)
add_text_box(slide, Inches(0.7), Inches(3.1), Inches(11), Inches(0.6),
    "kubectl: get pod, describe, logs, top  |  bash: curl, system commands  |  MCP: 확장 도구 (DB, HTTP 등)",
    10, False, GRAY_TEXT)

# Safety
safety_shape = add_shape(slide, Inches(0.5), Inches(4.1), Inches(12), Inches(1), RGBColor(0xfe, 0xf5, 0xe7), ORANGE)
add_text_box(slide, Inches(0.7), Inches(4.2), Inches(3), Inches(0.3), "안전 검증", 12, True, ORANGE)
add_text_box(slide, Inches(0.7), Inches(4.5), Inches(11), Inches(0.5),
    "O 명령어 분류  |  O 위험도 분석  |  ⚠ 위험 명령어 → 사용자 승인",
    10, False, GRAY_TEXT)

# Result
result_shape = add_shape(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(1.5), DARK_BG)
add_text_box(slide, Inches(0.7), Inches(5.4), Inches(11), Inches(0.3), "최종 답변", 12, True, GREEN)
add_text_box(slide, Inches(0.7), Inches(5.7), Inches(11), Inches(0.4),
    "redis-master의 메모리 제한이 256Mi로 설정되어 있으나 최근 트래픽 증가로 380Mi까지 사용 중입니다.",
    10, False, RGBColor(0xcc, 0xcc, 0xcc))
add_text_box(slide, Inches(0.7), Inches(6.1), Inches(11), Inches(0.4),
    "실행: kubectl set resources deployment/redis-master -c redis --limits=memory=512Mi",
    10, False, RGBColor(0xcc, 0xcc, 0xcc))

add_slide_number(slide, 7, TOTAL_SLIDES)

# Slide 8: Overview
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "주요 개선 작업", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

cards = [
    ("Web UI 고도화", "Lucide 아이콘 시스템 교체, TailAdmin 스타일 사이드바, 다크모드 지원, 한글 i18n 전체 적용으로 사용자 경험 대폭 향상", ACCENT_RED),
    ("CLI 신규 개발", "터미널 기반 REPL 인터페이스 구축, 한글 입력 지원, K8s 컨텍스트 인식 AI 명령어로 개발자 생산성 강화", BLUE),
    ("TUI 개선", "화면 깜빡임 3단계 개선, Diff 렌더링 구현, 웰컴 스크린 추가로 터미널 UI 품질 향상", GREEN),
    ("AI 모델 설정 방식 개선", "AI 모델 설정 UI 개선, openrouter 제공 모델 연동 기능 추가", PURPLE),
]

for i, (title, desc, color) in enumerate(cards):
    row, col = divmod(i, 2)
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 2.8)
    
    shape = add_shape(slide, x, y, Inches(6), Inches(2.5), WHITE)
    shape.line.color.rgb = color
    shape.line.width = Pt(3)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(10)
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = GRAY_TEXT

# Stats
stats = [("37", "Total Commits"), ("73", "Files Changed"), ("+10,939", "Lines Added"), ("-1,770", "Lines Removed")]
for i, (num, label) in enumerate(stats):
    x = Inches(1 + i * 3)
    add_text_box(slide, x, Inches(6.5), Inches(2.5), Inches(0.5), num, 28, True, ACCENT_RED, PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(7), Inches(2.5), Inches(0.3), label, 12, False, GRAY_TEXT, PP_ALIGN.CENTER)

add_slide_number(slide, 8, TOTAL_SLIDES)

# Slide 9: Web UI Features
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "Web UI 주요 업데이트", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

web_features = [
    ("Lucide 아이콘 시스템", "기존 아이콘 → Lucide 전체 교체\n상단 헤더 아이콘만 표시로 변경하여 사용성 개선", ACCENT_RED),
    ("다크 모드 지원", "헤더 토글 버튼 추가, Sun/Moon 아이콘\nOllama 미니멀리스트 테마 방식", BLUE),
    ("한글 i18n", "전체 UI 번역, 언어 설정 영어/한국어 제한\nAI 패널 번역", GREEN),
    ("Cluster Visualizer", "Pod 간 통신 실시간 애니메이션 효과\nK8S 전체 상태 시각화", PURPLE),
    ("AI 어시스턴트 개선", "세션 지우기 버튼, 도구 실행 정보 유지\nopenRouter 연동 추가", ORANGE),
    ("로그인 화면", "브랜딩 정리, 미니멀리스트 디자인", RGBColor(0xc0, 0x39, 0x2b)),
]

for i, (title, desc, color) in enumerate(web_features):
    row, col = divmod(i, 3)
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.3 + row * 3)
    
    shape = add_shape(slide, x, y, Inches(3.9), Inches(2.7), WHITE)
    shape.line.color.rgb = RGBColor(0xe0, 0xe0, 0xe0)
    shape.line.width = Pt(0.5)
    
    icon = add_shape(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5), color)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.9)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(8)
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY_TEXT

add_slide_number(slide, 9, TOTAL_SLIDES)

# Slides 10-13: Screenshots (placeholder)
for i in range(4):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT_GRAY
    
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "Web UI 개선 화면", 32, True, DARK_TEXT)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_RED
    shape.line.fill.background()
    
    shape = add_shape(slide, Inches(1), Inches(1.3), Inches(11.3), Inches(5.8), WHITE, RGBColor(0xe0, 0xe0, 0xe0))
    add_text_box(slide, Inches(4), Inches(3.5), Inches(5), Inches(1), f"screenshot-web{i+1}.png", 16, False, GRAY_TEXT, PP_ALIGN.CENTER)
    
    add_slide_number(slide, 10 + i, TOTAL_SLIDES)

# Slide 14: TUI Screenshot
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "TUI 개선 화면", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

shape = add_shape(slide, Inches(1), Inches(1.3), Inches(11.3), Inches(5.8), WHITE, RGBColor(0xe0, 0xe0, 0xe0))
add_text_box(slide, Inches(4), Inches(3.5), Inches(5), Inches(1), "screenshot-tui.png", 16, False, GRAY_TEXT, PP_ALIGN.CENTER)

add_slide_number(slide, 14, TOTAL_SLIDES)

# Slide 15: CLI Features
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "CLI 추가 (신규)", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

# Left panel
left_shape = add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.8), WHITE, RGBColor(0xe0, 0xe0, 0xe0))
tf = left_shape.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "주요 기능"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_TEXT

cli_features = [
    "NEW  REPL 구현 (1,554줄)",
    "NEW  한글(UTF-8) 입력 지원",
    "NEW  Left/Right 커서 네비게이션",
    "NEW  Ctrl+Enter 개행 입력",
    "NEW  도구 승인 프롬프트 (Y/A/n/q)",
    "NEW  K8s 컨텍스트 인식 AI",
    "NEW  Tab 자동 완성",
    "NEW  명령 히스토리 (Up/Down)",
]
for feat in cli_features:
    p2 = tf.add_paragraph()
    p2.text = feat
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_TEXT
    p2.space_before = Pt(6)

# Right panel - commands
right_shape = add_shape(slide, Inches(6.3), Inches(1.3), Inches(6.5), Inches(5.8), WHITE, RGBColor(0xe0, 0xe0, 0xe0))
tf = right_shape.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.2)
p = tf.paragraphs[0]
p.text = "명령어 리스트"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_TEXT

commands = [
    ":help  -  도움말 표시",
    ":ai <질문>  -  AI 어시스턴트에 질문",
    ":model [이름]  -  LLM 모델 조회/전환",
    ":namespace [이름]  -  네임스페이스 조회/설정",
    ":context [이름]  -  K8s 컨텍스트 조회/전환",
    ":mcp [list|tools]  -  MCP 상태/도구 목록",
    ":history  -  명령 히스토리 표시",
    ":clear  -  화면 지우기",
    ":version  -  버전 정보 표시",
    ":quit / :exit  -  CLI 종료",
]
for cmd in commands:
    p2 = tf.add_paragraph()
    p2.text = cmd
    p2.font.size = Pt(11)
    p2.font.color.rgb = DARK_TEXT
    p2.font.name = "Monaco"
    p2.space_before = Pt(4)

add_slide_number(slide, 15, TOTAL_SLIDES)

# Slide 16: CLI Screenshot
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "CLI 실제 구현 화면", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

shape = add_shape(slide, Inches(1), Inches(1.3), Inches(11.3), Inches(5.8), WHITE, RGBColor(0xe0, 0xe0, 0xe0))
add_text_box(slide, Inches(4), Inches(3.5), Inches(5), Inches(1), "screenshot-cli.png", 16, False, GRAY_TEXT, PP_ALIGN.CENTER)

add_slide_number(slide, 16, TOTAL_SLIDES)

# Slide 17: Technical
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "기술적 도전과 해결", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

tech_items = [
    ("한글 입력 처리", "CLI 터미널에서 다중바이트 UTF-8 문자의 정확한 입력과 커서 이동 처리 구현\nreadLine 함수 개선으로 한국어 사용자 지원 강화"),
    ("AI 도구 승인 플로우", "stdin 기반 세션 자동 승인(A all)\nsafety.PolicyEnforcer와 연동하여 안전한 도구 실행 보장"),
    ("K8s 컨텍스트 수집 방식 개선", "네임스페이스/리소스 현황 자동 조회\nYAML/Events/Logs 통합 제공으로 AI 분석 품질 향상"),
    ("Web UI 번역 시스템", "i18n.js 734줄 대폭 개선. AI 패널\nDecision 팝업, Cluster Overview 등 핵심 UI 한글 번역 완성"),
]

for i, (title, desc) in enumerate(tech_items):
    row, col = divmod(i, 2)
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 3)
    
    shape = add_shape(slide, x, y, Inches(6), Inches(2.7), WHITE)
    shape.line.color.rgb = RGBColor(0xe0, 0xe0, 0xe0)
    shape.line.width = Pt(0.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(10)
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY_TEXT

add_slide_number(slide, 17, TOTAL_SLIDES)

# Slide 18: Direction
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "K13D의 방향성 고민", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

directions = [
    ("타겟 유저", "Kubernetes를 운영하는 플랫폼 엔지니어, AI로 K8s 운영을 자동화하려는 DevOps 팀, CLI · TUI를 선호하는 개발자"),
    ("공략할 틈새 시장", "오픈소스 K8s AI 네이티브 전용 도구 특화, Air-gapped(폐쇄망)/온프레미스 보안 환경용 도구, k8s 자격증 대비 테스트 도구"),
    ("CLI 기능 확대 전략", "AI 기반 문제 진단 및 자동 수정, 반복 작업 자동화 스크립트 지원, MCP 기반 플러그인 생태계 구축"),
    ("실무 활용 예시 발굴", "장애 상황 AI 자동 진단 · 복구, 배포 전 영향도 분석 리포트, 인시던트 기반 사후 분석"),
]

for i, (title, desc) in enumerate(directions):
    row, col = divmod(i, 2)
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 3)
    
    shape = add_shape(slide, x, y, Inches(6), Inches(2.7), WHITE)
    shape.line.color.rgb = RGBColor(0xe0, 0xe0, 0xe0)
    shape.line.width = Pt(0.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.space_after = Pt(10)
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY_TEXT

add_slide_number(slide, 18, TOTAL_SLIDES)

# Slide 19: Summary
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = DARK_BG

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.5), "총 정리", 28, True, WHITE)

stats_summary = [
    ("14", "Web UI 커밋"),
    ("9", "CLI 커밋"),
    ("4", "TUI 커밋"),
    ("10", "버그 수정/문서화"),
]
for i, (num, label) in enumerate(stats_summary):
    row, col = divmod(i, 2)
    x = Inches(0.8 + col * 2.5)
    y = Inches(1.3 + row * 1.5)
    
    shape = add_shape(slide, x, y, Inches(2.2), Inches(1.2), RGBColor(0x25, 0x25, 0x40))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
    p2.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(7), Inches(0.5), Inches(5), Inches(0.5), "향후 계획", 28, True, WHITE)

phases = [
    ("Phase 1", "Web UI 성능 최적화 및 반응형 디자인 보완"),
    ("Phase 2", "사용 가능한 AI 모델과의 최적 연동성 강화"),
    ("Phase 3", "다양한 WebUI / CLI 기능 발굴 및 개선"),
    ("Phase 4", "기존 사용중인 도구에서 부착해서 사용 가능한 플러그인 사업 모델 개발"),
]
for i, (phase, desc) in enumerate(phases):
    y = Inches(1.3 + i * 1.2)
    
    shape = add_shape(slide, Inches(7), y, Inches(1.2), Inches(0.4), ACCENT_RED)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = phase
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    add_text_box(slide, Inches(8.4), y, Inches(4.5), Inches(0.4), desc, 13, False, RGBColor(0xcc, 0xcc, 0xcc))

add_text_box(slide, Inches(2), Inches(6.5), Inches(9), Inches(0.5),
    "k13d - Kubernetes Management, Reimagined", 16, False, RGBColor(0xcc, 0xcc, 0xcc), PP_ALIGN.CENTER)

add_slide_number(slide, 19, TOTAL_SLIDES)

# Slide 20: CLI Reference
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = LIGHT_GRAY

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.6), "참고용 CLI UI", 32, True, DARK_TEXT)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.95), Inches(0.8), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_RED
shape.line.fill.background()

refs = ["hermes.png", "ibmbob.png", "mimo.png", "opencode.png"]
for i, ref in enumerate(refs):
    row, col = divmod(i, 2)
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 3)
    
    shape = add_shape(slide, x, y, Inches(6), Inches(2.7), WHITE, RGBColor(0xe0, 0xe0, 0xe0))
    add_text_box(slide, x + Inches(1), y + Inches(1), Inches(4), Inches(0.5), ref, 14, False, GRAY_TEXT, PP_ALIGN.CENTER)

add_slide_number(slide, 20, TOTAL_SLIDES)

# Save
output_path = os.path.join(os.path.dirname(__file__), "k13d_리뉴얼_발표자료.pptx")
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
